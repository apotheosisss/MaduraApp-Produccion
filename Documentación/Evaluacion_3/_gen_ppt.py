# -*- coding: utf-8 -*-
"""Genera la presentación de defensa Eval 3 (.pptx) de MaduraApp."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_D = RGBColor(0x1B, 0x5E, 0x20)
AMBER = RGBColor(0xF9, 0xA8, 0x25)
DARK = RGBColor(0x1A, 0x1C, 0x19)
GRAY = RGBColor(0x55, 0x5B, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def bullets(slide, x, y, w, h, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2); r.font.color.rgb = color
        r.font.name = "Calibri"
        p.space_after = Pt(6)
    return tb

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def content_header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.1), GREEN)
    rect(slide, 0, Inches(1.1), SW, Pt(4), AMBER)
    textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
            title, size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        textbox(slide, Inches(0.6), Inches(0.0), Inches(12), Inches(0.4),
                kicker, size=11, color=RGBColor(0xCF,0xE8,0xCF))

def table(slide, x, y, w, h, headers, rows, col_w=None):
    nr, nc = len(rows) + 1, len(headers)
    gfx = slide.shapes.add_table(nr, nc, x, y, w, h)
    t = gfx.table
    if col_w:
        for i, cw in enumerate(col_w):
            t.columns[i].width = Inches(cw)
    for i, hh in enumerate(headers):
        c = t.cell(0, i); c.text = hh
        c.fill.solid(); c.fill.fore_color.rgb = GREEN
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        pr.runs[0].font.size = Pt(13); pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            pr = c.text_frame.paragraphs[0]
            run = pr.add_run(); run.text = str(val) if val else " "
            run.font.size = Pt(12); run.font.color.rgb = DARK
    return t

# ─────────── SLIDE 1 — PORTADA ───────────
s = add_slide()
rect(s, 0, 0, SW, SH, GREEN)
rect(s, 0, Inches(5.0), SW, Pt(5), AMBER)
textbox(s, Inches(1), Inches(2.1), Inches(11.3), Inches(1.2),
        "MaduraApp", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
        "Análisis de Madurez Agrícola mediante Visión Computacional",
        size=20, color=RGBColor(0xE8,0xF5,0xE9), align=PP_ALIGN.CENTER, italic=True)
textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
        "Defensa — Estado de Avance 3 · TPY1101", size=22, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.0), Inches(11.3), Inches(1.0),
        "Claudio Vicente Aro Kath · Sección 001D · Docente: José Ignacio Campos Arévalo",
        size=14, color=RGBColor(0xCF,0xE8,0xCF), align=PP_ALIGN.CENTER)
notes(s, "Buenos días/tardes. Soy Claudio Aro y presento la defensa del Estado de Avance 3 de "
         "MaduraApp, un sistema de análisis de madurez de frutas con visión computacional. En los "
         "próximos 15 minutos mostraré el plan de pruebas, su aplicación a los componentes del "
         "proyecto y las mejoras que derivaron de los resultados. Saludar con seguridad y mirar al jurado.")

# ─────────── SLIDE 2 — AGENDA ───────────
s = add_slide(); content_header(s, "Agenda")
bullets(s, Inches(1.2), Inches(1.6), Inches(11), Inches(5),
        [("Contexto general del proyecto", 0),
         ("Descripción del plan de pruebas", 0),
         ("Aplicación del plan y resultados", 0),
         ("Mejoras realizadas al producto", 0),
         ("Conclusión y trabajo pendiente", 0)], size=24)
notes(s, "Esta es la ruta de la presentación. Primero el contexto para situar el problema, luego el "
         "corazón de esta evaluación: el plan de pruebas, su aplicación y resultados, y las mejoras "
         "que hicimos a partir de esos resultados. Cierro con conclusión y lo que queda pendiente. "
         "Mencionar que el foco de la Eval 3 es calidad y seguridad.")

# ─────────── SLIDE 3 — CONTEXTO ───────────
s = add_slide(); content_header(s, "Contexto general del proyecto")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.2), Inches(5),
        [("Problema: 20–40% de pérdidas post-cosecha en frutas climatéricas (FAO/ODEPA)", 0),
         ("Causa raíz: falta de criterios objetivos de madurez", 0),
         ("Solución: app móvil + IA que clasifica madurez y recomienda acción", 0),
         ("4 frutas × 3 estados: Inmaduro / Óptimo / Sobre maduro", 1),
         ("Stack: Android (Kotlin) + FastAPI + YOLO26n", 0),
         ("Modelo: mAP@50 = 0.9229 (KPI ≥ 0.75)", 1)], size=18)
rect(s, Inches(8.3), Inches(1.7), Inches(4.3), Inches(4.2), LIGHT)
textbox(s, Inches(8.5), Inches(1.9), Inches(3.9), Inches(0.6), "Flujo del usuario", size=16, bold=True, color=GREEN_D)
bullets(s, Inches(8.5), Inches(2.5), Inches(3.9), Inches(3.2),
        [("Login / registro", 0), ("Elegir fruta", 0), ("Foto o galería", 0),
         ("Diagnóstico + semáforo", 0), ("Calificar resultado", 0), ("Historial (offline)", 0)], size=15)
notes(s, "El problema real: en Chile se pierde entre 20 y 40% de la fruta post-cosecha, en parte "
         "porque no hay una forma objetiva y accesible de saber el punto óptimo de consumo. "
         "MaduraApp resuelve eso con una foto: el usuario elige la fruta, toma la foto, y recibe el "
         "estado de madurez con un semáforo y una recomendación. El backend en FastAPI corre el "
         "modelo YOLO26n, que alcanza 0.92 de mAP, muy por encima del KPI de 0.75. Si preguntan por "
         "qué YOLO: es detección en tiempo real, ligero (5 MB), corre en CPU.")

# ─────────── SLIDE 4 — PLAN DE PRUEBAS ───────────
s = add_slide(); content_header(s, "Descripción del plan de pruebas", "Criterio 1 · IL3.1")
bullets(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5),
        [("57 casos automatizados (38 backend + 19 Android)", 0),
         ("Clasificados en 4 tipos:", 0),
         ("Validación: ¿el sistema correcto? (requisitos)", 1),
         ("Verificación: ¿correctamente? (calidad/RNF)", 1),
         ("Seguridad: ¿protegido? (OWASP)", 1),
         ("Operacional: ¿funciona en su entorno?", 1),
         ("Cada caso: funcionalidad · acción · esperado · obtenido", 0)], size=17)
table(s, Inches(7.0), Inches(1.7), Inches(5.7), Inches(2.5),
      ["Tipo", "Casos", "Estado"],
      [["Estado de fruta (backend)", "21", "✅"],
       ["Validación backend", "17", "✅"],
       ["Validación Android", "19", "✅"]],
      col_w=[3.1, 1.3, 1.3])
textbox(s, Inches(7.0), Inches(4.4), Inches(5.7), Inches(0.6),
        "Base de pruebas: SQLite in-memory + fixtures + mocks", size=14, italic=True, color=GRAY)
notes(s, "El plan de pruebas está alineado a la problemática: cada componente crítico tiene "
         "pruebas. Lo organicé en cuatro tipos —y aquí está la diferencia que el jurado valora—: "
         "validación responde '¿construyo el sistema correcto?' (cumple los requisitos); "
         "verificación, '¿lo construyo correctamente?' (calidad, performance); seguridad, con "
         "criterios OWASP; y operacional, que funcione en su entorno. La base de datos de pruebas es "
         "SQLite en memoria, que se crea y destruye en cada corrida para aislar los tests, con "
         "fixtures que generan un usuario y un token reales. Si preguntan por qué in-memory: "
         "aislamiento, velocidad y no tocar datos reales.")

# ─────────── SLIDE 5 — APLICACIÓN Y RESULTADOS ───────────
s = add_slide(); content_header(s, "Aplicación del plan y resultados", "Criterio 2 · IL3.1")
table(s, Inches(0.7), Inches(1.6), Inches(7.2), Inches(3),
      ["Suite", "Framework", "Pruebas", "Resultado"],
      [["Backend FastAPI", "pytest", "38", "✅ 38/38"],
       ["Android datos", "MockK", "9", "✅ 9/9"],
       ["Android ViewModels", "JUnit", "10", "✅ 10/10"],
       ["TOTAL", "", "57", "✅ 57/57"]],
      col_w=[2.6, 1.7, 1.4, 1.5])
bullets(s, Inches(8.2), Inches(1.6), Inches(4.4), Inches(4.5),
        [("57/57 en verde", 0),
         ("21 verifican estado de fruta", 1),
         ("Ejecución en CI en cada push", 0),
         ("Rendimiento: ~200 ms inferencia", 0),
         ("0 errores hasta 50 concurrentes", 1),
         ("mAP@50 = 0.9229 · 5.2 MB", 0)], size=16)
notes(s, "Apliqué el plan y este es el resultado: 57 de 57 pruebas en verde, ejecutadas hoy. El "
         "backend con pytest, el Android con MockK y JUnit. Además agregué integración continua: un "
         "workflow que corre la suite en cada push, así ningún cambio rompe las pruebas sin que nos "
         "enteremos. En calidad: el modelo cumple el KPI con 0.92, pesa 5.2 MB y el APK compila. "
         "Si preguntan qué detecta cada test: los de validación verifican los caminos felices y de "
         "error de cada endpoint; los de Android, que el cache offline y los estados de la UI sean "
         "correctos. Demostrar dominio: puedo abrir el reporte si lo piden.")

# ─────────── SLIDE 6 — MEJORAS ───────────
s = add_slide(); content_header(s, "Mejoras realizadas al producto", "Criterio 3 · IL3.2")
textbox(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
        "15 mejoras derivadas de las pruebas y de una auditoría OWASP, por estándar de calidad:",
        size=16, color=GRAY)
table(s, Inches(0.7), Inches(1.9), Inches(12), Inches(3.5),
      ["Estándar", "Mejoras destacadas"],
      [["Seguridad", "Secreto JWT en prod, CORS acotado, política de contraseña, HTTPS forzado, token cifrado AES-256"],
       ["Usabilidad", "Rediseño Material 3: íconos vectoriales, modo oscuro, feedback táctil, tipografía consistente"],
       ["Corrección", "Suite y build a verde tras integración; integración continua (CI)"],
       ["Completitud", "Autenticación + feedback integrados y probados end-to-end"],
       ["Pertinencia", "Diagnóstico y recomendaciones ajustados al dominio agrícola"]],
      col_w=[2.4, 9.6])
notes(s, "Lo más importante de esta evaluación: las mejoras NO son arbitrarias, salen de los "
         "resultados de las pruebas y de la auditoría OWASP. Las mapeé a los cinco estándares que "
         "pide la rúbrica. En seguridad: encontré que el secreto del token estaba hardcodeado, que "
         "el tráfico iba en claro y que el token se guardaba sin cifrar; lo corregí todo. En "
         "usabilidad rediseñé la interfaz con Material 3. En corrección, dejé la suite y el build en "
         "verde. Cada mejora tiene un commit que la respalda. Pregunta probable: '¿por qué no es "
         "cifrado de extremo a extremo como WhatsApp?' Respuesta: porque el servidor necesita leer "
         "la imagen para inferir; sería deshonesto llamarlo E2E, así que uso 'cifrado en tránsito'.")

# ─────────── SLIDE 7 — SEGURIDAD (destacado) ───────────
s = add_slide(); content_header(s, "Foco: seguridad de datos personales", "OWASP")
bullets(s, Inches(0.7), Inches(1.6), Inches(7.2), Inches(5),
        [("Contraseñas hasheadas con bcrypt (nunca en texto plano)", 0),
         ("Endpoints protegidos con JWT (A01)", 0),
         ("Cifrado en tránsito: HTTPS forzado (A02)", 0),
         ("Cifrado en reposo: token AES-256 en el dispositivo (M9)", 0),
         ("Política de contraseña: 8+ con letra y número (A07)", 0),
         ("Secreto JWT obligatorio en producción (A05)", 0)], size=18)
rect(s, Inches(8.3), Inches(1.7), Inches(4.3), Inches(2.2), LIGHT)
textbox(s, Inches(8.5), Inches(1.85), Inches(3.9), Inches(2.0),
        "Mensaje honesto al usuario:\n«Tus datos viajan cifrados (HTTPS) y tu contraseña se guarda protegida»",
        size=15, color=GREEN_D, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Profundizo en seguridad porque protege datos personales del usuario. Apliqué OWASP punto "
         "por punto: contraseñas con bcrypt, endpoints con JWT, HTTPS forzado para el tránsito, y el "
         "token cifrado con AES-256 en el dispositivo para el reposo. Y algo de ética: en vez de "
         "mentir diciendo 'extremo a extremo', le muestro al usuario un mensaje veraz: cifrado en "
         "tránsito. Si preguntan qué es OWASP: es el estándar de referencia de seguridad en "
         "aplicaciones; A01 es control de acceso roto, A02 fallos criptográficos, A07 fallos de "
         "autenticación.")

# ─────────── SLIDE 8 — CONCLUSIÓN ───────────
s = add_slide(); content_header(s, "Conclusión y trabajo pendiente")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.2), Inches(5),
        [("Producto funcional, probado y seguro", 0),
         ("57/57 pruebas en verde en todos los componentes", 1),
         ("15 mejoras trazables a los 5 estándares de calidad", 1),
         ("Datos personales protegidos (cifrado + hashing + JWT)", 1),
         ("Pendiente honesto:", 0),
         ("Desplegar backend en AWS Lab del docente", 1),
         ("Pruebas E2E automatizadas (Espresso)", 1),
         ("Aprobación del plan en esta defensa", 1)], size=17)
rect(s, Inches(8.3), Inches(1.7), Inches(4.3), Inches(3.6), GREEN)
textbox(s, Inches(8.5), Inches(2.0), Inches(3.9), Inches(3.0),
        "57/57\npruebas en verde\n\nmAP@50 = 0.9229\n\n15 mejoras\nverificadas",
        size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Para cerrar: MaduraApp es un producto funcional, probado y seguro. 57 de 57 pruebas en "
         "verde, 15 mejoras trazables a los estándares de calidad, y datos personales protegidos. "
         "Soy honesto con lo que falta: desplegar el backend en el laboratorio AWS, automatizar las "
         "pruebas end-to-end, y obtener su aprobación del plan en esta defensa. Gracias, quedo "
         "atento a sus preguntas. Cerrar con seguridad y agradecer.")

# ─────────── SLIDE 9 — CIERRE ───────────
s = add_slide()
rect(s, 0, 0, SW, SH, GREEN)
rect(s, 0, Inches(4.3), SW, Pt(5), AMBER)
textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
        "¿Preguntas?", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.8),
        "MaduraApp · Claudio Vicente Aro Kath · TPY1101 — Sección 001D",
        size=16, color=RGBColor(0xCF,0xE8,0xCF), align=PP_ALIGN.CENTER)
notes(s, "Gracias. Quedo atento a sus preguntas. Recordar: respondo con calma, si no sé algo lo "
         "reconozco y explico cómo lo averiguaría. Las preguntas pueden ser sobre el plan de "
         "pruebas, la base de datos de pruebas, OWASP, o por qué tomé cada decisión técnica.")

out = os.path.join(os.path.dirname(__file__), "MaduraApp_Presentacion_Evaluacion3.pptx")
prs.save(out)
print("Guardado:", out, "·", len(prs.slides._sldIdLst), "slides")
